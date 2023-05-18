import os
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from google.oauth2 import service_account
from googleapiclient.discovery import build

# Function to extract mathematical content from .docx file


def extract_math_content(docx_file):
    doc = Document(docx_file)
    math_content = [paragraph.text for paragraph in doc.paragraphs if paragraph._element.xml.endswith(
        "w:instrText></w:instrText>")]
    return math_content

# Function to create a slide with an equation in PowerPoint presentation


def create_powerpoint_slide(presentation, equation):
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    textbox = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(8), Inches(4))
    textbox.text = equation

# Main function


def main():
    docx_file = '/home/nikhil/Documents/Internship/Sample(doc2ppt).docx'
    powerpoint_file = '/home/nikhil/Documents/Internship/powerpoint1.pptx'

    math_content = extract_math_content(docx_file)

    # Create PowerPoint presentation
    presentation = Presentation()
    for equation in math_content:
        create_powerpoint_slide(presentation, equation)
    presentation.save(powerpoint_file)


if __name__ == '__main__':
    main()
