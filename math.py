import os
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from google.oauth2 import service_accountls
from googleapiclient.discovery import build

# Function to extract mathematical content from .docx file


def extract_math_content(docx_file):
    doc = Document(docx_file)
    math_content = []

    for paragraph in doc.paragraphs:
        if paragraph._element.xml.endswith("w:instrText></w:instrText>"):
            math_content.append(paragraph.text)

    return math_content

# Function to create a PowerPoint slide with an equation


def create_powerpoint_slide(presentation, equation):
    # Use the layout for a content slide
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)

    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(4)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = equation

# Function to create a Google Slides slide with an equation


def create_google_slide(service, presentation_id, equation):
    body = {
        'requests': [{
            'createSlide': {
                'slideLayoutReference': {
                    'predefinedLayout': 'BLANK'
                }
            }
        }]
    }
    response = service.presentations().batchUpdate(
        presentationId=presentation_id, body=body).execute()

    slide_id = response['replies'][0]['createSlide']['objectId']
    element_id = response['replies'][0]['createSlide']['objectId']

    requests = [
        {
            'createSheetsChart': {
                'objectId': element_id,
                'spreadsheetId': '',
                'chartId': '',
                'linkingMode': 'LINKED',
                'elementProperties': {
                    'pageObjectId': slide_id,
                    'size': {
                        'width': {
                            'magnitude': 300,
                            'unit': 'PT'
                        },
                        'height': {
                            'magnitude': 200,
                            'unit': 'PT'
                        }
                    },
                    'transform': {
                        'scaleX': 1,
                        'scaleY': 1,
                        'translateX': 100,
                        'translateY': 100,
                        'unit': 'PT'
                    }
                }
            }
        }
    ]
    body = {'requests': requests}
    service.presentations().batchUpdate(
        presentationId=presentation_id, body=body).execute()

# Main function


def main():
    docx_file = '/home/nikhil/Documents/Internship/Sample(doc2ppt).docx'
    powerpoint_file = '/home/nikhil/Documents/Internship/powerpoint.pptx'
    google_slides_credentials = '/home/nikhil/Documents/Internship/credentials.json'
    google_slides_file = '/home/nikhil/Documents/Internship/google_slides'

    math_content = extract_math_content(docx_file)

    # Create PowerPoint presentation
    presentation = Presentation()
    for equation in math_content:
        create_powerpoint_slide(presentation, equation)
    presentation.save(powerpoint_file)

    # Create Google Slides presentation
    credentials = service_account.Credentials.from_service_account_file(
        google_slides_credentials, scopes=['https://www.googleapis.com/auth/presentations'])
    service = build('slides', 'v1', credentials=credentials)
    presentation = service.presentations().create().execute()
    presentation_id = presentation['presentationId']
    for equation in math_content:
        create_google_slide(service, presentation_id, equation)
    service.presentations().export(presentationId=presentation_id,
                                   mimeType='application/pdf').execute()
    os.rename(f'{presentation_id}.pdf', f'{google_slides_file}.pdf')


if __name__ == '__main__':
    main()
